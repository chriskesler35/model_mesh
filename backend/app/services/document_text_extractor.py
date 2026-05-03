"""Document text extraction helpers for chat file attachments."""

from __future__ import annotations

from io import BytesIO
import csv
from typing import Any


MAX_EXTRACTED_CHARS = 60_000


class ExtractionError(Exception):
    """Raised when a file cannot be extracted."""


def _decode_text(raw_bytes: bytes) -> str:
    for encoding in ("utf-8", "utf-8-sig", "utf-16", "cp1252", "latin-1"):
        try:
            return raw_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw_bytes.decode("utf-8", errors="replace")


def _truncate_text(text: str) -> tuple[str, bool]:
    if len(text) <= MAX_EXTRACTED_CHARS:
        return text, False
    return text[:MAX_EXTRACTED_CHARS], True


def _extract_txt_like(raw_bytes: bytes) -> str:
    return _decode_text(raw_bytes)


def _extract_csv(raw_bytes: bytes) -> str:
    text = _decode_text(raw_bytes)
    rows: list[str] = []
    reader = csv.reader(text.splitlines())
    for row in reader:
        rows.append(" | ".join(col.strip() for col in row))
    return "\n".join(rows)


def _extract_pdf(raw_bytes: bytes) -> str:
    try:
        from pypdf import PdfReader
    except Exception as exc:  # pragma: no cover - optional dependency guard
        raise ExtractionError("PDF support is not installed (missing pypdf).") from exc

    reader = PdfReader(BytesIO(raw_bytes))
    parts: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        page_text = (page.extract_text() or "").strip()
        if page_text:
            parts.append(f"[Page {i}]\n{page_text}")
    return "\n\n".join(parts)


def _extract_docx(raw_bytes: bytes) -> str:
    try:
        from docx import Document
    except Exception as exc:  # pragma: no cover
        raise ExtractionError("DOCX support is not installed (missing python-docx).") from exc

    document = Document(BytesIO(raw_bytes))
    parts: list[str] = []
    for paragraph in document.paragraphs:
        value = paragraph.text.strip()
        if value:
            parts.append(value)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def _extract_xlsx(raw_bytes: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except Exception as exc:  # pragma: no cover
        raise ExtractionError("XLSX support is not installed (missing openpyxl).") from exc

    workbook = load_workbook(filename=BytesIO(raw_bytes), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                lines.append(" | ".join(cells))
        lines.append("")
    return "\n".join(lines).strip()


def _extract_pptx(raw_bytes: bytes) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise ExtractionError("PPTX support is not installed (missing python-pptx).") from exc

    presentation = Presentation(BytesIO(raw_bytes))
    lines: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"[Slide {index}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                value = shape.text.strip()
                if value:
                    lines.append(value)
        lines.append("")
    return "\n".join(lines).strip()


def extract_document_text(*, filename: str, content_type: str | None, raw_bytes: bytes, max_chars: int = MAX_EXTRACTED_CHARS) -> dict[str, Any]:
    lower_name = (filename or "").lower()
    extension = lower_name.rsplit(".", 1)[-1] if "." in lower_name else ""

    extractor_by_ext = {
        "txt": _extract_txt_like,
        "md": _extract_txt_like,
        "py": _extract_txt_like,
        "js": _extract_txt_like,
        "ts": _extract_txt_like,
        "tsx": _extract_txt_like,
        "jsx": _extract_txt_like,
        "json": _extract_txt_like,
        "yaml": _extract_txt_like,
        "yml": _extract_txt_like,
        "csv": _extract_csv,
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "xlsx": _extract_xlsx,
        "pptx": _extract_pptx,
    }

    extractor = extractor_by_ext.get(extension)
    if extractor is None and (content_type or "").startswith("text/"):
        extractor = _extract_txt_like

    if extractor is None:
        raise ExtractionError(
            f"Unsupported file type: '{extension or 'unknown'}'. Supported: pdf, docx, xlsx, pptx, txt, csv, md, json, and common code/text files."
        )

    extracted = extractor(raw_bytes)
    cleaned = extracted.strip()
    if not cleaned:
        cleaned = "(No extractable text found in this file.)"
    capped = min(max(1000, max_chars), MAX_EXTRACTED_CHARS)
    if len(cleaned) <= capped:
        truncated_text, truncated = cleaned, False
    else:
        truncated_text, truncated = cleaned[:capped], True

    return {
        "file_name": filename,
        "mime_type": content_type or "application/octet-stream",
        "extension": extension,
        "extracted_text": truncated_text,
        "extracted_chars": len(truncated_text),
        "truncated": truncated,
    }
